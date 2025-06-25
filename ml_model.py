import os
import pandas as pd
import numpy as np
from sklearn.ensemble import RandomForestClassifier
from sklearn.preprocessing import StandardScaler, LabelEncoder
import joblib
from loguru import logger
from utils import get_file_extension, calculate_sha256
import docx
import pptx
import PyPDF2
import pefile
import magic
import hashlib
from datetime import datetime


class FileAnalyzer:
    def __init__(self):
        self.model_path = os.path.join('models', 'ids_rf_model.joblib')
        self.scaler_path = os.path.join('models', 'ids_scaler.joblib')
        self.encoders_path = os.path.join('models', 'ids_encoders.joblib')
        self.model = None
        self.scaler = None
        self.label_encoders = None
        # Define protocol mapping for file extensions
        self.protocol_map = {
            '.docx': 1,
            '.pptx': 2,
            '.pdf': 3,
            '.exe': 4,
            '.dll': 5,
            '.bat': 6,
            '.ps1': 7,
            '.vbs': 8,
            '.js': 9,
            '.py': 10,
            '.txt': 11,
            '.csv': 12
        }
        self._load_or_train_model()

    def _load_or_train_model(self):
        if all(os.path.exists(p) for p in [self.model_path, self.scaler_path, self.encoders_path]):
            try:
                self.model = joblib.load(self.model_path)
                self.scaler = joblib.load(self.scaler_path)
                self.label_encoders = joblib.load(self.encoders_path)
                logger.info('Loaded ML model, scaler, and encoders.')
            except Exception as e:
                logger.error(f'Error loading model: {e}')
                self._train_model()
        else:
            self._train_model()

    def _train_model(self):
        try:
            df = pd.read_csv(os.path.join('data', 'cybersecurity_attacks.csv'))
            categorical = [c for c in ['Protocol', 'Attack Type'] if c in df.columns]
            self.label_encoders = {}
            for col in categorical:
                le = LabelEncoder()
                df[col] = le.fit_transform(df[col].astype(str))
                self.label_encoders[col] = le
            features = [c for c in ['Source Port', 'Destination Port', 'Packet Length', 'Protocol'] if c in df.columns]
            X = df[features].fillna(0)
            y = (df['Attack Type'].notna().astype(int) if 'Attack Type' in df.columns else np.zeros(len(df)))
            self.scaler = StandardScaler()
            X_scaled = self.scaler.fit_transform(X)
            self.model = RandomForestClassifier(n_estimators=100, random_state=42)
            self.model.fit(X_scaled, y)
            os.makedirs('models', exist_ok=True)
            joblib.dump(self.model, self.model_path)
            joblib.dump(self.scaler, self.scaler_path)
            joblib.dump(self.label_encoders, self.encoders_path)
            logger.info('Trained and saved new ML model.')
        except Exception as e:
            logger.error(f'Error training model: {e}')

    def analyze_file(self, filepath):
        try:
            ext = get_file_extension(filepath)
            size = os.path.getsize(filepath)
            sha256 = calculate_sha256(filepath)
            
            # Extract file-specific features
            features = self.extract_file_features(filepath)
            if features is None:
                return {
                    'is_malicious': False,
                    'summary': 'Could not analyze file contents',
                    'sha256': sha256
                }
            
            # Analyze based on file type and content
            is_malicious = False
            risk_factors = []
            file_type = magic.from_file(filepath, mime=True)
            
            # Detailed analysis based on file type
            if ext.lower() in ['.exe', '.dll']:
                try:
                    pe = pefile.PE(filepath)
                    # Mark all .exe and .dll files as malicious as per user request
                    is_malicious = True
                    risk_factors.append("Marked as malicious by policy for executable files")
                    summary = f"Executable file with {len(pe.sections)} sections"
                    if hasattr(pe, 'DIRECTORY_ENTRY_IMPORT'):
                        imports = len(pe.DIRECTORY_ENTRY_IMPORT)
                        summary += f", {imports} imports"
                except Exception as e:
                    summary = f"Executable file (Error analyzing PE: {str(e)})"
                    is_malicious = True
                    risk_factors.append("Could not analyze executable structure")
            
            elif ext.lower() in ['.bat', '.ps1', '.vbs', '.js']:
                try:
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    if size == 0:
                        risk_factors.append("Empty script file")
                        is_malicious = True
                    elif size < 10:
                        risk_factors.append("Suspiciously small script")
                        is_malicious = True
                    lines = content.count('\n') + 1
                    summary = f"Script file with {lines} lines, {len(content)} characters"
                    if len(content.strip()) == 0:
                        risk_factors.append("Script contains only whitespace")
                        is_malicious = True
                except Exception as e:
                    summary = f"Script file (Error reading content: {str(e)})"
                    is_malicious = True
                    risk_factors.append("Could not read script content")
            
            elif ext.lower() in ['.pdf']:
                try:
                    with open(filepath, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        pages = len(reader.pages)
                        summary = f"PDF document with {pages} pages"
                        if pages == 0:
                            risk_factors.append("Empty PDF document")
                            is_malicious = True
                except Exception as e:
                    summary = f"PDF document (Error analyzing: {str(e)})"
                    is_malicious = True
                    risk_factors.append("Could not analyze PDF structure")
            
            elif ext.lower() in ['.docx']:
                try:
                    doc = docx.Document(filepath)
                    text = ' '.join([p.text for p in doc.paragraphs])
                    summary = f"Word document with {len(doc.paragraphs)} paragraphs, {len(text)} characters"
                    if len(text.strip()) == 0:
                        risk_factors.append("Empty document")
                        is_malicious = True
                except Exception as e:
                    summary = f"Word document (Error analyzing: {str(e)})"
                    is_malicious = True
                    risk_factors.append("Could not analyze document content")
            
            elif ext.lower() in ['.txt', '.csv']:
                try:
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    lines = content.count('\n') + 1
                    summary = f"Text file with {lines} lines, {len(content)} characters"
                    if len(content.strip()) == 0:
                        risk_factors.append("Empty text file")
                        is_malicious = True
                except Exception as e:
                    summary = f"Text file (Error reading: {str(e)})"
                    is_malicious = True
                    risk_factors.append("Could not read file content")
            
            else:
                summary = f"File of type {file_type}, size {size} bytes"
                if size == 0:
                    risk_factors.append("Empty file")
                    is_malicious = True
            
            # Add risk factors to summary if any found
            if risk_factors:
                summary += f" - Risk Factors: {', '.join(risk_factors)}"
            
            return {
                'is_malicious': is_malicious,
                'summary': summary,
                'sha256': sha256
            }
            
        except Exception as e:
            logger.error(f"Error analyzing file {filepath}: {str(e)}")
            return {
                'is_malicious': False,
                'summary': f'Error analyzing file: {str(e)}',
                'sha256': calculate_sha256(filepath) if os.path.exists(filepath) else 'N/A'
            }

    def _summarize_file(self, filepath, ext, size):
        try:
            if ext == '.docx':
                doc = docx.Document(filepath)
                text = ' '.join([p.text for p in doc.paragraphs])
                return f'Word document containing {len(text)} characters of text'
            elif ext == '.pptx':
                prs = pptx.Presentation(filepath)
                slides = len(prs.slides)
                return f'PowerPoint presentation with {slides} slides'
            elif ext == '.pdf':
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    pages = len(reader.pages)
                return f'PDF document with {pages} pages'
            elif ext == '.exe':
                pe = pefile.PE(filepath)
                return f'Executable file with {len(pe.sections)} sections'
            elif ext == '.bat':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return f'Batch script containing {len(content)} characters'
            elif ext == '.txt':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return f'Text file containing {len(content)} characters'
            elif ext == '.py':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return f'Python script containing {len(content)} characters'
            elif ext == '.js':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return f'JavaScript file containing {len(content)} characters'
            elif ext == '.ps1':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return f'PowerShell script containing {len(content)} characters'
            elif ext == '.vbs':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return f'VBScript file containing {len(content)} characters'
            else:
                return f'File of size {size} bytes'
        except Exception as e:
            return f'Could not read file contents: {str(e)}'

    def extract_file_features(self, file_path):
        """Extract features from a file for analysis."""
        try:
            features = {}
            
            # Basic file features
            file_size = os.path.getsize(file_path)
            features['file_size'] = file_size
            
            # File type using python-magic
            file_type = magic.from_file(file_path, mime=True)
            features['file_type'] = file_type
            
            # Calculate file hash
            with open(file_path, 'rb') as f:
                file_hash = hashlib.sha256(f.read()).hexdigest()
            features['file_hash'] = file_hash
            
            # Calculate file entropy
            features['entropy'] = self._calculate_entropy(file_hash)
            
            # Extract features based on file type
            if file_type in ['application/x-dosexec', 'application/x-msdos-program']:
                try:
                    pe = pefile.PE(file_path)
                    features.update({
                        'number_of_sections': len(pe.sections),
                        'has_resources': int(hasattr(pe, 'DIRECTORY_ENTRY_RESOURCE')),
                        'is_dll': int(pe.is_dll()),
                        'is_exe': int(pe.is_exe()),
                        'timestamp': pe.FILE_HEADER.TimeDateStamp
                    })
                except:
                    pass
            
            return features
            
        except Exception as e:
            logger.error(f"Error extracting features from {file_path}: {e}")
            return None
    
    def _calculate_entropy(self, data):
        """Calculate Shannon entropy of the data."""
        if not data:
            return 0
        entropy = 0
        for x in range(256):
            p_x = data.count(chr(x)) / len(data)
            if p_x > 0:
                entropy += -p_x * np.log2(p_x)
        return entropy
    
    def get_file_summary(self, file_path):
        """Generate a summary of the file based on its type."""
        try:
            # Use python-magic to get file type
            file_type = magic.from_file(file_path, mime=True)
            summary = {
                'file_type': file_type,
                'file_size': os.path.getsize(file_path),
                'last_modified': datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat(),
                'md5_hash': self._calculate_md5(file_path)
            }
            
            # Add type-specific analysis based on MIME type
            if file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                summary.update(self._analyze_docx(file_path))
            elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                summary.update(self._analyze_pptx(file_path))
            elif file_type == 'application/pdf':
                summary.update(self._analyze_pdf(file_path))
            elif file_type in ['application/x-dosexec', 'application/x-msdos-program']:
                summary.update(self._analyze_pe(file_path))
            else:
                # For other file types, just return basic info
                summary['summary'] = f'File type: {file_type}, Size: {summary["file_size"]} bytes'
            
            return summary
            
        except Exception as e:
            logger.error(f"Error generating file summary for {file_path}: {e}")
            return {
                'error': str(e),
                'file_type': 'unknown',
                'file_size': os.path.getsize(file_path) if os.path.exists(file_path) else 0,
                'summary': f'Error during analysis: {str(e)}'
            }
    
    def _calculate_md5(self, file_path):
        """Calculate MD5 hash of a file."""
        try:
            with open(file_path, 'rb') as f:
                return hashlib.md5(f.read()).hexdigest()
        except:
            return None
    
    def _analyze_pe(self, file_path):
        """Analyze PE (Portable Executable) files."""
        try:
            pe = pefile.PE(file_path)
            return {
                'is_dll': pe.is_dll(),
                'is_exe': pe.is_exe(),
                'number_of_sections': len(pe.sections),
                'number_of_imports': len(pe.DIRECTORY_ENTRY_IMPORT) if hasattr(pe, 'DIRECTORY_ENTRY_IMPORT') else 0,
                'has_resources': hasattr(pe, 'DIRECTORY_ENTRY_RESOURCE'),
                'compilation_time': datetime.fromtimestamp(pe.FILE_HEADER.TimeDateStamp).isoformat()
            }
        except Exception as e:
            logger.error(f"Error analyzing PE file {file_path}: {e}")
            return {}
    
    def _analyze_docx(self, file_path):
        """Analyze DOCX files."""
        from docx import Document
        try:
            doc = Document(file_path)
            return {
                'page_count': len(doc.paragraphs),
                'has_images': any(shape for shape in doc.inline_shapes),
                'metadata': {
                    'author': doc.core_properties.author,
                    'created': doc.core_properties.created,
                    'modified': doc.core_properties.modified,
                    'last_modified_by': doc.core_properties.last_modified_by
                }
            }
        except Exception as e:
            logger.error(f"Error analyzing DOCX {file_path}: {e}")
            return {}
    
    def _analyze_pptx(self, file_path):
        """Analyze PPTX files."""
        from pptx import Presentation
        try:
            prs = Presentation(file_path)
            return {
                'slide_count': len(prs.slides),
                'has_images': any(slide.shapes for slide in prs.slides),
                'metadata': {
                    'author': prs.core_properties.author,
                    'created': prs.core_properties.created,
                    'modified': prs.core_properties.modified,
                    'last_modified_by': prs.core_properties.last_modified_by
                }
            }
        except Exception as e:
            logger.error(f"Error analyzing PPTX {file_path}: {e}")
            return {}
    
    def _analyze_pdf(self, file_path):
        """Analyze PDF files."""
        from PyPDF2 import PdfReader
        try:
            reader = PdfReader(file_path)
            return {
                'page_count': len(reader.pages),
                'is_encrypted': reader.is_encrypted,
                'metadata': reader.metadata,
                'has_attachments': any('/Filespec' in str(page) for page in reader.pages)
            }
        except Exception as e:
            logger.error(f"Error analyzing PDF {file_path}: {e}")
            return {} 